"""Block builders: every slide-deck.schema.json block type as editable shapes."""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import charts
from .theme import Theme, est_lines, set_run_font

IN = 914400
NA_TEXT = {"zh": "暂无数据", "en": "N/A"}

TRANSFORM_SCALE = {"raw": 1, "percent": 1, "multiple": 1, "wan": 1e4, "yi": 1e8,
                   "thousand": 1e3, "million": 1e6, "billion": 1e9}
TRANSFORM_SUFFIX = {"percent": "%", "multiple": "x"}


@dataclass
class Resolver:
    """Resolves metric/claim/assumption references from the research package."""

    package: dict
    metrics: dict = field(init=False)
    claims: dict = field(init=False)
    assumptions: dict = field(init=False)

    def __post_init__(self) -> None:
        self.metrics = {m["metric_id"]: m for m in self.package["metrics"]}
        self.claims = {c["claim_id"]: c for c in self.package["claims"]}
        self.assumptions = {a["assumption_id"]: a
                            for a in self.package["valuation"]["assumptions"]}

    def value(self, display_number: dict) -> float:
        metric = self.metrics[display_number["metric_id"]]
        scale = TRANSFORM_SCALE[display_number["display_transform"]]
        return round(metric["value"] / scale, display_number.get("decimals", 1))

    def display(self, display_number: dict, lang: str) -> str:
        v = self.value(display_number)
        d = display_number.get("decimals", 1)
        text = f"{v:,.{d}f}"
        suffix = TRANSFORM_SUFFIX.get(display_number["display_transform"], "")
        return text + suffix


def sign_color(theme: Theme, value: float):
    if value > 0:
        return theme.positive
    if value < 0:
        return theme.negative
    return theme.muted


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(int(0.06 * IN))
    tf.margin_top = tf.margin_bottom = Emu(int(0.02 * IN))
    return box, tf


def _para(tf, first: bool):
    return tf.paragraphs[0] if first else tf.add_paragraph()


# ---------------------------------------------------------------- kpi_cards --

def build_kpi_cards(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                    x: Emu, y: Emu, w: Emu) -> Emu:
    cards = block["kpi_cards"]
    gap = Emu(int(0.18 * IN))
    card_w = Emu(int((w - gap * (len(cards) - 1)) / len(cards)))
    card_h = Emu(int(1.28 * IN))
    for i, card in enumerate(cards):
        cx = Emu(x + i * (card_w + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, card_w, card_h)
        shape.adjustments[0] = 0.08
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.tint
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(int(0.12 * IN))
        tf.margin_top = Emu(int(0.10 * IN))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = card["label"]
        set_run_font(run, theme, lang=lang, size=11, color=theme.primary, bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        num = card["number"]
        run2 = p2.add_run()
        run2.text = res.display(num, lang)
        color = theme.dark
        if num.get("show_sign_colour"):
            color = sign_color(theme, res.value(num))
        set_run_font(run2, theme, lang=lang, size=theme.size_pt("kpi_number", lang) - 6,
                     color=color, bold=True)
        if card.get("period_caption"):
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.LEFT
            run3 = p3.add_run()
            run3.text = card["period_caption"]
            set_run_font(run3, theme, lang=lang, size=9, color=theme.muted)
    return Emu(y + card_h + Emu(int(0.22 * IN)))


# ------------------------------------------------------------------ bullets --

def build_bullets(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                  x: Emu, y: Emu, w: Emu) -> Emu:
    marker = "・ " if lang == "zh" else "– "
    size = theme.size_pt("body", lang)
    box, tf = _textbox(slide, x, y, w, Emu(int(0.4 * IN)))
    total_lines = 0
    for i, bullet in enumerate(block["bullets"]):
        p = _para(tf, i == 0)
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = marker + bullet["text"]
        set_run_font(run, theme, lang=lang, size=size)
        tag = {"fact": ("事实", "Fact"), "inference": ("推断", "Inference"),
               "opinion": ("观点", "Opinion")}[bullet["claim_type"]]
        tag_run = p.add_run()
        tag_run.text = f"  [{tag[0] if lang == 'zh' else tag[1]}]"
        set_run_font(tag_run, theme, lang=lang, size=9, color=theme.light, bold=True)
        total_lines += est_lines(run.text, size, w / IN)
    h = Emu(int(total_lines * size * theme.line_height(lang) / 72 * IN + 0.24 * IN))
    box.height = h
    return Emu(y + h + Emu(int(0.12 * IN)))


# -------------------------------------------------------------------- chart --

def build_chart(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                x: Emu, y: Emu, w: Emu, h: Emu) -> Emu:
    spec = block["chart"]
    categories = spec.get("x_labels") or [""]
    series: list[tuple[str, list[float | None]]] = []
    for s in spec["series"]:
        values = []
        for mid in s["metric_ids"]:
            scale = TRANSFORM_SCALE[s.get("display_transform", "raw")]
            values.append(round(res.metrics[mid]["value"] / scale, 2))
        series.append((s["name"], values))
    ctype = spec["chart_type"]
    if ctype == "waterfall":
        comps = list(zip(categories, series[0][1]))
        categories, series = charts.waterfall_series(comps)
        charts.add_chart(slide, theme, lang, chart_type=ctype, categories=categories,
                         series=series, x=x, y=y, w=w, h=h,
                         signal_colors=[theme.primary, theme.positive, theme.negative])
    else:
        charts.add_chart(slide, theme, lang, chart_type=ctype, categories=categories,
                         series=series, x=x, y=y, w=w, h=h)
    title_y = Emu(y + h)
    caption = spec["title"]
    if spec.get("unit_caption"):
        caption += f" · {spec['unit_caption']}"
    if ctype == "waterfall":
        caption += " · " + ("紫=总额 绿=增加 红=减少" if lang == "zh"
                            else "purple=total, green=increase, red=decrease")
    caption += " · " + ("来源 " if lang == "zh" else "Source ") + ", ".join(spec["source_ids"])
    charts.add_chart_caption(slide, theme, lang, text=caption, x=x, y=title_y, w=w)
    return Emu(title_y + Emu(int(0.30 * IN)))


# -------------------------------------------------------------------- table --

def build_table(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                x: Emu, y: Emu, w: Emu) -> Emu:
    spec = block["table"]
    cols = spec["columns"]
    rows = spec["rows"]
    n_rows = len(rows) + 1
    row_h = Emu(int(0.32 * IN))
    shape = slide.shapes.add_table(n_rows, len(cols), x, y, w, Emu(int(row_h * n_rows)))
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    label_w = Emu(int(w * 0.24))
    data_w = Emu(int((w - label_w) / (len(cols) - 1)))
    table.columns[0].width = label_w
    for ci in range(1, len(cols)):
        table.columns[ci].width = data_w
    for ci, name in enumerate(cols):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.primary
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = name
        set_run_font(run, theme, lang=lang, size=theme.size_pt("table_body", lang),
                     color=theme.inverse, bold=True)
    for ri, row in enumerate(rows, start=1):
        label_cell = table.cell(ri, 0)
        label_cell.fill.solid()
        label_cell.fill.fore_color.rgb = theme.mid
        p = label_cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = row["label"]
        set_run_font(run, theme, lang=lang, size=theme.size_pt("table_body", lang),
                     color=theme.inverse)
        for ci, cell_spec in enumerate(row["cells"], start=1):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.surface if ri % 2 else theme.tint
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            if cell_spec is None:
                run.text = NA_TEXT[lang]
                set_run_font(run, theme, lang=lang, size=theme.size_pt("table_body", lang),
                             color=theme.muted)
            else:
                run.text = res.display(cell_spec, lang)
                color = theme.body_text
                if cell_spec.get("show_sign_colour"):
                    color = sign_color(theme, res.value(cell_spec))
                set_run_font(run, theme, lang=lang, size=theme.size_pt("table_body", lang),
                             color=color)
    cap_y = Emu(y + row_h * n_rows + Emu(int(0.04 * IN)))
    caption = spec["unit_caption"] + " · " + ("来源 " if lang == "zh" else "Source ") + \
        ", ".join(spec["source_ids"])
    charts.add_chart_caption(slide, theme, lang, text=caption, x=x, y=cap_y, w=w)
    return Emu(cap_y + Emu(int(0.30 * IN)))


# -------------------------------------------------------- comparison_cards --

def build_comparison_cards(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                           x: Emu, y: Emu, w: Emu) -> Emu:
    cards = block["comparison_cards"]
    gap = Emu(int(0.2 * IN))
    card_w = Emu(int((w - gap * (len(cards) - 1)) / len(cards)))
    card_h = Emu(int(3.3 * IN))
    for i, card in enumerate(cards):
        cx = Emu(x + i * (card_w + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, card_w, card_h)
        shape.adjustments[0] = 0.05
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.surface
        shape.line.color.rgb = theme.light
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(int(0.14 * IN))
        tf.margin_top = Emu(int(0.12 * IN))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = card["entity_name"]
        set_run_font(run, theme, lang=lang, size=14, color=theme.primary, bold=True,
                     heading=True)
        if card.get("limited_data"):
            lim = p.add_run()
            lim.text = "  " + ("（资料有限）" if lang == "zh" else "(limited data)")
            set_run_font(lim, theme, lang=lang, size=9, color=theme.muted)
        for numspec in card.get("numbers", []):
            p2 = tf.add_paragraph()
            run_l = p2.add_run()
            run_l.text = numspec["label"] + "  "
            set_run_font(run_l, theme, lang=lang, size=10, color=theme.muted)
            run_v = p2.add_run()
            run_v.text = res.display(numspec["number"], lang)
            set_run_font(run_v, theme, lang=lang, size=13, color=theme.dark, bold=True)
        for b in card["bullets"]:
            p3 = tf.add_paragraph()
            p3.space_before = Pt(4)
            run_b = p3.add_run()
            run_b.text = ("・ " if lang == "zh" else "– ") + b["text"]
            set_run_font(run_b, theme, lang=lang, size=11)
    return Emu(y + card_h + Emu(int(0.15 * IN)))


# --------------------------------------------------------------------- flow --

def build_flow(slide, theme: Theme, res: Resolver, lang: str, block: dict,
               x: Emu, y: Emu, w: Emu) -> Emu:
    nodes = block["flow"]
    gap = Emu(int(0.28 * IN))
    node_w = Emu(int((w - gap * (len(nodes) - 1)) / len(nodes)))
    node_h = Emu(int(1.55 * IN))
    for i, node in enumerate(nodes):
        nx = Emu(x + i * (node_w + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, nx, y, node_w, node_h)
        shape.adjustments[0] = 0.18
        shape.fill.solid()
        highlighted = node.get("highlight")
        shape.fill.fore_color.rgb = theme.primary if highlighted else theme.light
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(int(0.18 * IN))
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = node["node_name"]
        set_run_font(run, theme, lang=lang, size=13,
                     color=theme.inverse if highlighted else theme.dark, bold=True)
        if node.get("description"):
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = node["description"]
            set_run_font(run2, theme, lang=lang, size=10,
                         color=theme.inverse if highlighted else theme.body_text)
        if node.get("number"):
            p3 = tf.add_paragraph()
            run3 = p3.add_run()
            run3.text = res.display(node["number"], lang)
            set_run_font(run3, theme, lang=lang, size=12,
                         color=theme.inverse if highlighted else theme.dark, bold=True)
    return Emu(y + node_h + Emu(int(0.2 * IN)))


# ------------------------------------------------------------------ timeline --

def build_timeline(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                   x: Emu, y: Emu, w: Emu) -> Emu:
    events = block["timeline"]
    line_y = Emu(y + Emu(int(0.55 * IN)))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, line_y, w, Emu(int(0.02 * IN)))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.primary
    line.line.fill.background()
    line.shadow.inherit = False
    step = Emu(int(w / len(events)))
    dot_d = Emu(int(0.14 * IN))
    for i, ev in enumerate(events):
        cx = Emu(x + i * step + Emu(int(0.1 * IN)))
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(cx), Emu(line_y - dot_d // 2 + Emu(int(0.01 * IN))),
            dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme.dark if not ev.get("is_future") else theme.light
        dot.line.color.rgb = theme.dark
        dot.line.width = Pt(0.75)
        dot.shadow.inherit = False
        # date above
        box, tf = _textbox(slide, Emu(cx - Emu(int(0.2 * IN))), y,
                           Emu(step - Emu(int(0.1 * IN))), Emu(int(0.35 * IN)))
        run = tf.paragraphs[0].add_run()
        run.text = ev["date"]
        set_run_font(run, theme, lang=lang, size=10, color=theme.primary, bold=True)
        # label + description below
        box2, tf2 = _textbox(slide, Emu(cx - Emu(int(0.2 * IN))),
                             Emu(line_y + Emu(int(0.18 * IN))),
                             Emu(step - Emu(int(0.1 * IN))), Emu(int(1.5 * IN)))
        run2 = tf2.paragraphs[0].add_run()
        run2.text = ev["label"]
        set_run_font(run2, theme, lang=lang, size=11, color=theme.dark, bold=True)
        if ev.get("description"):
            p = tf2.add_paragraph()
            run3 = p.add_run()
            run3.text = ev["description"]
            set_run_font(run3, theme, lang=lang, size=10)
    return Emu(line_y + Emu(int(1.9 * IN)))


# -------------------------------------------------------------- football ----

def build_football_field(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                         x: Emu, y: Emu, w: Emu, h: Emu) -> Emu:
    spec = block["football_field"]
    bars = [(b["label"], res.value(b["low"]), res.value(b["high"])) for b in spec["bars"]]
    cats, series = charts.football_series(bars)
    chart_w = Emu(int(w * 0.62))
    lows = [b[1] for b in bars] + [res.value(spec["current_price"])]
    highs = [b[2] for b in bars] + [res.value(spec["current_price"])]
    charts.add_chart(slide, theme, lang, chart_type="football_field",
                     categories=cats, series=series, x=x, y=y, w=chart_w, h=h,
                     signal_colors=[theme.primary],
                     axis_min=float(int(min(lows) * 0.9)),
                     axis_max=float(int(max(highs) * 1.08) + 1))
    # current/target price lines rendered as labelled markers next to the chart
    info_x = Emu(x + chart_w + Emu(int(0.3 * IN)))
    info_w = Emu(w - chart_w - Emu(int(0.3 * IN)))
    box, tf = _textbox(slide, info_x, y, info_w, h)
    pairs = [("现价" if lang == "zh" else "Current price", spec["current_price"], theme.dark)]
    if spec.get("target_price"):
        pairs.append(("目标价" if lang == "zh" else "Target price",
                      spec["target_price"], theme.primary))
    first = True
    for label, dn, color in pairs:
        p = _para(tf, first)
        first = False
        run = p.add_run()
        run.text = f"{label}  "
        set_run_font(run, theme, lang=lang, size=11, color=theme.muted)
        run2 = p.add_run()
        run2.text = res.display(dn, lang)
        set_run_font(run2, theme, lang=lang, size=16, color=color, bold=True)
    for line in spec.get("assumption_lines", []):
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = ("・ " if lang == "zh" else "– ") + line["text"]
        set_run_font(run, theme, lang=lang, size=10)
    cap_y = Emu(y + h + Emu(int(0.02 * IN)))
    caption = ("估值区间 · 来源 " if lang == "zh" else "Valuation range · Source ") + \
        ", ".join(spec["source_ids"])
    charts.add_chart_caption(slide, theme, lang, text=caption, x=x, y=cap_y, w=w)
    return Emu(cap_y + Emu(int(0.3 * IN)))


# --------------------------------------------------------- paired_columns ---

def _parse_chip(chip: str) -> tuple[int, int] | None:
    """Parse machine-readable severity/likelihood chips like 'S:high L:medium'."""
    levels = {"low": 0, "medium": 1, "high": 2, "低": 0, "中": 1, "高": 2}
    s_val = l_val = None
    for part in chip.replace("，", " ").replace(",", " ").split():
        if ":" in part:
            key, _, val = part.partition(":")
            val = levels.get(val.strip().lower())
            if val is None:
                continue
            if key.strip().upper().startswith("S") or key.strip() in ("严重性",):
                s_val = val
            elif key.strip().upper().startswith("L") or key.strip() in ("概率", "可能性"):
                l_val = val
    if s_val is None or l_val is None:
        return None
    return s_val, l_val


def build_paired_columns(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                         x: Emu, y: Emu, w: Emu, *, with_matrix: bool = True) -> Emu:
    spec = block["paired_columns"]
    matrix_positions = [(_parse_chip(item.get("chip", "")), item)
                        for item in spec["right_items"]]
    has_matrix = with_matrix and any(pos for pos, _ in matrix_positions)
    col_w = Emu(int((w - Emu(int(0.4 * IN))) * (0.31 if has_matrix else 0.5)))
    matrix_w = Emu(w - col_w * 2 - Emu(int(0.8 * IN))) if has_matrix else Emu(0)

    def render_column(items, title, cx, accent):
        box, tf = _textbox(slide, cx, y, col_w, Emu(int(4.4 * IN)))
        run = tf.paragraphs[0].add_run()
        run.text = title
        set_run_font(run, theme, lang=lang, size=15, color=accent, bold=True, heading=True)
        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            run1 = p.add_run()
            run1.text = ("・ " if lang == "zh" else "– ") + item["title"]
            set_run_font(run1, theme, lang=lang, size=12, color=theme.dark, bold=True)
            if item.get("chip"):
                chip_run = p.add_run()
                parsed = _parse_chip(item["chip"])
                if parsed:
                    lv = ["低", "中", "高"] if lang == "zh" else ["Low", "Med", "High"]
                    chip_text = (f"严重性{lv[parsed[0]]}·概率{lv[parsed[1]]}" if lang == "zh"
                                 else f"Severity {lv[parsed[0]]} · Likelihood {lv[parsed[1]]}")
                else:
                    chip_text = item["chip"]
                chip_run.text = f"  [{chip_text}]"
                set_run_font(chip_run, theme, lang=lang, size=9, color=accent, bold=True)
            if item.get("description"):
                p2 = tf.add_paragraph()
                run2 = p2.add_run()
                run2.text = "   " + item["description"]
                set_run_font(run2, theme, lang=lang, size=10.5)

    render_column(spec["left_items"], spec["left_title"], x, theme.positive)
    render_column(spec["right_items"], spec["right_title"],
                  Emu(x + col_w + Emu(int(0.4 * IN))), theme.negative)

    if has_matrix:
        mx = Emu(x + col_w * 2 + Emu(int(0.8 * IN)))
        cell = Emu(int(matrix_w / 3))
        m_h = Emu(cell * 3)
        my = Emu(y + Emu(int(0.42 * IN)))
        shades = [theme.tint, theme.light, theme.mid]
        for row in range(3):        # severity: top = high
            for col in range(3):    # likelihood: left = low
                sev = 2 - row
                shade = shades[max(sev, col) if max(sev, col) < 3 else 2]
                sq = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(mx + col * cell), Emu(my + row * cell),
                    Emu(cell - Emu(int(0.03 * IN))), Emu(cell - Emu(int(0.03 * IN))))
                sq.fill.solid()
                sq.fill.fore_color.rgb = shade
                sq.line.fill.background()
                sq.shadow.inherit = False
        for idx, (pos, item) in enumerate([(p, i) for p, i in matrix_positions if p]):
            sev, lik = pos
            dot_d = Emu(int(0.26 * IN))
            dx = Emu(mx + lik * cell + cell // 2 - dot_d // 2)
            dy = Emu(my + (2 - sev) * cell + cell // 2 - dot_d // 2)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dx, dy, dot_d, dot_d)
            dot.fill.solid()
            dot.fill.fore_color.rgb = theme.dark
            dot.line.color.rgb = theme.surface
            dot.line.width = Pt(1.25)
            dot.shadow.inherit = False
            tfd = dot.text_frame
            run = tfd.paragraphs[0].add_run()
            run.text = f"R{idx + 1}"
            set_run_font(run, theme, lang=lang, size=8, color=theme.inverse, bold=True)
        # axis captions
        cap, tfc = _textbox(slide, mx, Emu(my + m_h + Emu(int(0.03 * IN))),
                            Emu(cell * 3), Emu(int(0.25 * IN)))
        runc = tfc.paragraphs[0].add_run()
        runc.text = ("可能性 →（纵轴：严重性 ↑）" if lang == "zh"
                     else "Likelihood → (vertical: severity ↑)")
        set_run_font(runc, theme, lang=lang, size=9, color=theme.muted)
    return Emu(y + Emu(int(4.6 * IN)))


# -------------------------------------------------------------- source_table -

def build_source_table(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                       x: Emu, y: Emu, w: Emu) -> Emu:
    rows = block["source_table"]
    headers = (["编号", "标题", "发布方", "发布日期", "获取日期"] if lang == "zh"
               else ["ID", "Title", "Publisher", "Published", "Retrieved"])
    widths = [0.10, 0.42, 0.22, 0.13, 0.13]
    shape = slide.shapes.add_table(len(rows) + 1, 5, x, y, w,
                                   Emu(int(0.28 * IN) * (len(rows) + 1)))
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    for ci, header in enumerate(headers):
        table.columns[ci].width = Emu(int(w * widths[ci]))
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.primary
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = header
        set_run_font(run, theme, lang=lang, size=10, color=theme.inverse, bold=True)
    for ri, src in enumerate(rows, start=1):
        values = [src["source_id"], src.get("title", ""), src.get("publisher", ""),
                  src.get("published_date", "—"), src.get("retrieved_at", "")]
        for ci, value in enumerate(values):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.surface if ri % 2 else theme.tint
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = str(value)
            set_run_font(run, theme, lang=lang, size=9,
                         color=theme.body_text)
    return Emu(y + Emu(int(0.28 * IN) * (len(rows) + 1)) + Emu(int(0.12 * IN)))


# ---------------------------------------------------------------- text_panel -

def build_text_panel(slide, theme: Theme, res: Resolver, lang: str, block: dict,
                     x: Emu, y: Emu, w: Emu) -> Emu:
    spec = block["text_panel"]
    style = spec.get("style", "note")
    size = theme.size_pt("body", lang) if style == "conclusion" else 10
    fill = theme.tint if style == "conclusion" else None
    lines = est_lines(spec["text"], size, w / IN)
    h = Emu(int((lines * size * theme.line_height(lang) / 72 + 0.28) * IN))
    if fill is not None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.adjustments[0] = 0.05
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(int(0.16 * IN))
    else:
        _, tf = _textbox(slide, x, y, w, h)
    run = tf.paragraphs[0].add_run()
    run.text = spec["text"]
    color = theme.dark if style == "conclusion" else theme.muted
    set_run_font(run, theme, lang=lang, size=size, color=color,
                 bold=style == "conclusion")
    return Emu(y + h + Emu(int(0.14 * IN)))
