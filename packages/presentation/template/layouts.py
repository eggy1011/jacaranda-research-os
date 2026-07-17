"""Slide compositors: header/footer/logo chrome plus L01-L11 block dispatch."""

from __future__ import annotations

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import blocks
from .blocks import Resolver
from .theme import Theme, set_run_font

import tempfile

from PIL import Image

IN = 914400
ASSETS = Path(__file__).resolve().parents[3] / "assets" / "brand"
LOGO_FULL = ASSETS / "jacaranda-logo-full.png"
LOGO_SHIELD = ASSETS / "jacaranda-shield.png"

_RESIZE_CACHE: dict[tuple[str, int], str] = {}


def brand_image(path: Path, height_in: float, dpi: int = 220) -> str:
    """Downscaled working copy of a brand master so decks stay small; master untouched."""
    px = int(height_in * dpi)
    key = (str(path), px)
    if key not in _RESIZE_CACHE:
        img = Image.open(path)
        ratio = px / img.height
        img = img.resize((max(1, int(img.width * ratio)), px), Image.LANCZOS)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp.name)
        _RESIZE_CACHE[key] = tmp.name
    return _RESIZE_CACHE[key]

# full lockup aspect (w/h) measured from the master asset
FULL_ASPECT = 3375 / 4016
SHIELD_ASPECT = 2727 / 2970


def _footer(slide, theme: Theme, lang: str, deck: dict, slide_spec: dict) -> None:
    footer = slide_spec["footer"]
    y = Emu(theme.page_h - theme.footer_h - Emu(int(0.08 * IN)))
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, theme.content_left, y,
                                  theme.content_width, Emu(int(0.008 * IN)))
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.light
    rule.line.fill.background()
    rule.shadow.inherit = False
    ty = Emu(y + Emu(int(0.05 * IN)))
    parts = []
    if footer["show_page_number"]:
        parts.append((f"第 {slide_spec['slide_no']} 页" if lang == "zh"
                      else f"Page {slide_spec['slide_no']}"))
    parts.append(("数据截至 " if lang == "zh" else "Data as of ") + footer["data_as_of"])
    if footer["source_ids"]:
        parts.append(("来源 " if lang == "zh" else "Source ") + ", ".join(footer["source_ids"]))
    if footer.get("show_disclaimer_ref", True):
        parts.append("仅供学习研究，不构成投资建议" if lang == "zh"
                     else "For educational research only; not investment advice")
    box = slide.shapes.add_textbox(theme.content_left, ty, theme.content_width,
                                   Emu(int(0.25 * IN)))
    tf = box.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text = "   ·   ".join(parts)
    set_run_font(run, theme, lang=lang, size=theme.size_pt("footnote", lang),
                 color=theme.muted)


def _header(slide, theme: Theme, lang: str, slide_spec: dict) -> None:
    title_box = slide.shapes.add_textbox(theme.content_left, Emu(int(0.42 * IN)),
                                         Emu(int(theme.content_width * 0.82)),
                                         Emu(int(0.55 * IN)))
    tf = title_box.text_frame
    tf.word_wrap = False
    run = tf.paragraphs[0].add_run()
    run.text = slide_spec["title"]
    set_run_font(run, theme, lang=lang, size=theme.size_pt("slide_title", lang),
                 color=theme.dark, bold=True, heading=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, theme.content_left,
                                  Emu(int(1.05 * IN)), theme.content_width,
                                  Emu(int(0.018 * IN)))
    rule.fill.solid()
    rule.fill.fore_color.rgb = theme.primary
    rule.line.fill.background()
    rule.shadow.inherit = False
    # shield top-right (shield-only fallback: the supplied master is a vertical lockup,
    # unreadable at 0.45in height — documented in TEMPLATE_GUIDE.md)
    logo_h = Inches(theme.tokens["layout"]["logo"]["content_slides"]["height_in"])
    logo_w = Emu(int(logo_h * SHIELD_ASPECT))
    slide.shapes.add_picture(brand_image(LOGO_SHIELD, logo_h / IN),
                             Emu(theme.page_w - theme.margin["right"] - logo_w),
                             Emu(int(0.32 * IN)), height=logo_h)


def _content_geometry(theme: Theme):
    return (theme.content_left, theme.content_top, theme.content_width,
            Emu(theme.content_bottom - theme.content_top))


BLOCK_BUILDERS = {
    "kpi_cards": blocks.build_kpi_cards,
    "bullets": blocks.build_bullets,
    "table": blocks.build_table,
    "comparison_cards": blocks.build_comparison_cards,
    "flow": blocks.build_flow,
    "timeline": blocks.build_timeline,
    "paired_columns": blocks.build_paired_columns,
    "source_table": blocks.build_source_table,
    "text_panel": blocks.build_text_panel,
}


def build_content_slide(slide, theme: Theme, res: Resolver, lang: str,
                        deck: dict, slide_spec: dict) -> None:
    _header(slide, theme, lang, slide_spec)
    _footer(slide, theme, lang, deck, slide_spec)
    x, y, w, h = _content_geometry(theme)
    cursor = y
    n_charts = sum(1 for b in slide_spec["blocks"] if b["block_type"] == "chart")
    for block in slide_spec["blocks"]:
        btype = block["block_type"]
        if btype == "chart":
            remaining = Emu(theme.content_bottom - cursor - Emu(int(0.35 * IN)))
            chart_h = Emu(int(min(remaining, Emu(int(3.6 * IN))) /
                              (1 if n_charts == 1 else 1)))
            if len(slide_spec["blocks"]) > 1:
                chart_h = Emu(int(min(chart_h, Emu(int(3.0 * IN)))))
            cursor = blocks.build_chart(slide, theme, res, lang, block,
                                        x, cursor, w, chart_h)
        elif btype == "football_field":
            remaining = Emu(theme.content_bottom - cursor - Emu(int(0.4 * IN)))
            cursor = blocks.build_football_field(slide, theme, res, lang, block,
                                                 x, cursor, w,
                                                 Emu(int(min(remaining,
                                                             Emu(int(3.4 * IN))))))
        else:
            cursor = BLOCK_BUILDERS[btype](slide, theme, res, lang, block, x, cursor, w)


def build_cover(slide, theme: Theme, res: Resolver, lang: str, deck: dict,
                slide_spec: dict) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, theme.page_w, theme.page_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.dark
    bg.line.fill.background()
    bg.shadow.inherit = False
    # white safe-area card behind the transparent-cutout lockup
    logo_h = Inches(2.4)
    logo_w = Emu(int(logo_h * FULL_ASPECT))
    card_pad = Emu(int(0.35 * IN))
    card_w = Emu(logo_w + card_pad * 2)
    card_h = Emu(logo_h + card_pad * 2)
    card_x = Emu(int((theme.page_w - card_w) / 2))
    card_y = Emu(int(0.55 * IN))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y,
                                  card_w, card_h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = theme.surface
    card.line.fill.background()
    card.shadow.inherit = False
    slide.shapes.add_picture(brand_image(LOGO_FULL, logo_h / IN), Emu(card_x + card_pad),
                             Emu(card_y + card_pad), height=logo_h)
    ty = Emu(card_y + card_h + Emu(int(0.35 * IN)))
    title_box = slide.shapes.add_textbox(Inches(1.2), ty,
                                         Emu(theme.page_w - Inches(2.4)),
                                         Emu(int(0.9 * IN)))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_spec["title"]
    set_run_font(run, theme, lang=lang, size=theme.size_pt("cover_title", lang) - 6,
                 color=theme.inverse, bold=True, heading=True)
    if slide_spec.get("kicker"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = slide_spec["kicker"]
        set_run_font(run2, theme, lang=lang, size=theme.size_pt("cover_subtitle", lang) - 4,
                     color=theme.light)
    meta = slide_spec["blocks"][0]["cover_meta"]
    lines = [meta["company_line"]]
    if meta.get("rating_line"):
        lines.append(meta["rating_line"])
    lines.append("   ·   ".join(filter(None, [meta.get("date_line"),
                                              meta.get("edition_line"),
                                              meta.get("prepared_by")])))
    my = Emu(theme.page_h - Emu(int(1.5 * IN)))
    mbox = slide.shapes.add_textbox(Inches(1.2), my, Emu(theme.page_w - Inches(2.4)),
                                    Emu(int(1.1 * IN)))
    mtf = mbox.text_frame
    mtf.word_wrap = True
    for i, line in enumerate(lines):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        set_run_font(run, theme, lang=lang, size=13 if i < 2 else 10,
                     color=theme.inverse if i < 2 else theme.light,
                     bold=i == 1)


def build_divider(slide, theme: Theme, res: Resolver, lang: str, deck: dict,
                  slide_spec: dict) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, theme.page_w, theme.page_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.background
    bg.line.fill.background()
    bg.shadow.inherit = False
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(int(2.9 * IN)),
                                  theme.page_w, Emu(int(1.7 * IN)))
    band.fill.solid()
    band.fill.fore_color.rgb = theme.primary
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(1.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_spec["title"]
    set_run_font(run, theme, lang=lang, size=30, color=theme.inverse, bold=True,
                 heading=True)
    if slide_spec.get("kicker"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = slide_spec["kicker"]
        set_run_font(run2, theme, lang=lang, size=14, color=theme.light)
    logo_h = Inches(theme.tokens["layout"]["logo"]["content_slides"]["height_in"])
    logo_w = Emu(int(logo_h * SHIELD_ASPECT))
    slide.shapes.add_picture(brand_image(LOGO_SHIELD, logo_h / IN),
                             Emu(theme.page_w - theme.margin["right"] - logo_w),
                             Emu(int(0.32 * IN)), height=logo_h)
    _footer(slide, theme, lang, deck, slide_spec)


LAYOUT_BUILDERS = {
    "L01_cover": build_cover,
    "L02_section_divider": build_divider,
}


def build_slide(prs, theme: Theme, res: Resolver, lang: str, deck: dict,
                slide_spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    layout = slide_spec["layout"]
    builder = LAYOUT_BUILDERS.get(layout, build_content_slide)
    builder(slide, theme, res, lang, deck, slide_spec)
    return slide
