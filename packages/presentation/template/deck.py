"""Deck assembly: validated slide-deck JSON + research package -> editable PPTX.

Emits the PresentationProvider-contract overflow/QA report: every entry carries
slide_no, a traceable block identifier, the action taken (or failure action) and a
reason — consumable by downstream retry/drop logic.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from .blocks import Resolver
from .theme import Theme, load_theme

IN = 914400
SUPPORTED_EDITIONS = {"zh-CN": "zh", "en-AU": "en"}


class TemplateRenderError(ValueError):
    """Typed render failure; `code` is machine-readable for the provider contract."""

    def __init__(self, code: str, detail: str = "") -> None:
        self.code = code
        self.detail = detail
        super().__init__(f"{code}: {detail}" if detail else code)


class UnsupportedEditionError(TemplateRenderError):
    pass


@dataclass
class BuildContext:
    """Collects per-block placement spans and overflow-policy events during a build."""

    theme: Theme
    events: list[dict] = field(default_factory=list)
    block_spans: dict[int, list[tuple[int, str, int, int, int]]] = field(default_factory=dict)

    def record_span(self, slide_no: int, index: int, btype: str, priority: int,
                    y0: int, y1: int) -> None:
        self.block_spans.setdefault(slide_no, []).append((index, btype, priority, y0, y1))

    def record_event(self, *, code: str, slide_no: int, layout: str, block: str,
                     priority: int, action_taken: str, retryable: bool, reason: str) -> None:
        self.events.append({
            "code": code, "slide_no": slide_no, "layout": layout, "block": block,
            "priority": priority, "action_taken": action_taken,
            "retryable": retryable, "reason": reason,
        })

    def block_id_at(self, slide_no: int, top: int, height: int) -> tuple[str, int]:
        """Map a shape's vertical span back to the block that produced it."""
        centre = top + height // 2
        for index, btype, priority, y0, y1 in self.block_spans.get(slide_no, []):
            if y0 - int(0.05 * IN) <= centre <= y1 + int(0.05 * IN):
                return f"blocks[{index}]/{btype}", priority
        return "slide_chrome", 1


def _edition_lang(edition: str) -> str:
    if edition not in SUPPORTED_EDITIONS:
        raise UnsupportedEditionError(
            "unsupported_edition",
            f"edition {edition!r} is not renderable by this template system "
            f"(supported: {sorted(SUPPORTED_EDITIONS)}); bilingual-summary is out of "
            "scope per Issue #24",
        )
    return SUPPORTED_EDITIONS[edition]


def build_deck(deck_json: dict, package_json: dict, out_path: Path,
               theme: Theme | None = None) -> dict:
    """Build the PPTX and return the machine-readable QA/overflow report."""
    from .layouts import build_slide  # local import to avoid cycle

    theme = theme or load_theme()
    lang = _edition_lang(deck_json["edition"])
    res = Resolver(package_json)
    ctx = BuildContext(theme=theme)
    prs = Presentation()
    prs.slide_width = theme.page_w
    prs.slide_height = theme.page_h
    for slide_spec in deck_json["slides"]:
        build_slide(prs, theme, res, lang, deck_json, slide_spec, ctx)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return qa_check(out_path, deck_json, theme, ctx)


def qa_check(pptx_path: Path, deck_json: dict, theme: Theme | None = None,
             ctx: BuildContext | None = None) -> dict:
    """Geometry + placeholder QA on the produced file (reopened from disk)."""
    theme = theme or load_theme()
    ctx = ctx or BuildContext(theme=theme)
    prs = Presentation(str(pptx_path))
    issues: list[dict] = list(ctx.events)
    page_w, page_h = prs.slide_width, prs.slide_height

    def walk_shapes(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == 6:  # group
                yield from walk_shapes(shape.shapes)

    for idx, slide in enumerate(prs.slides, start=1):
        spec = deck_json["slides"][idx - 1]
        for shape in walk_shapes(slide.shapes):
            left = shape.left if shape.left is not None else 0
            top = shape.top if shape.top is not None else 0
            width = shape.width or 0
            height = shape.height or 0
            if left < 0 or top < 0 or left + width > page_w + Emu(int(0.02 * IN)) \
                    or top + height > page_h + Emu(int(0.02 * IN)):
                block_id, priority = ctx.block_id_at(idx, top, height)
                issues.append({
                    "code": "shape_out_of_bounds", "slide_no": idx,
                    "layout": spec["layout"], "block": block_id, "priority": priority,
                    "action_taken": "none",
                    "retryable": True,
                    "reason": f"shape {shape.shape_id} at ({left/IN:.2f},{top/IN:.2f})in "
                              f"size ({width/IN:.2f}x{height/IN:.2f})in exceeds the page; "
                              "recompress or drop this block",
                })
            if shape.has_text_frame:
                text = shape.text_frame.text
                # NB: bare "XXX" is not scanned — fictional tickers use it (600XXX)
                for token in ("{{", "}}", "TODO", "PLACEHOLDER", "lorem"):
                    if token in text:
                        block_id, priority = ctx.block_id_at(idx, top, height)
                        issues.append({
                            "code": "unresolved_placeholder", "slide_no": idx,
                            "layout": spec["layout"], "block": block_id,
                            "priority": priority, "action_taken": "none",
                            "retryable": False,
                            "reason": f"template token {token!r} left in rendered text",
                        })
    return {
        "pptx": pptx_path.name,
        "edition": deck_json["edition"],
        "slide_count": len(prs.slides),
        "page_size_in": [round(page_w / IN, 3), round(page_h / IN, 3)],
        "issues": issues,
        "status": "pass" if not issues else "fail",
    }


def build_from_files(deck_path: Path, package_path: Path, out_path: Path) -> dict:
    deck_json = json.loads(deck_path.read_text(encoding="utf-8"))
    package_json = json.loads(package_path.read_text(encoding="utf-8"))
    return build_deck(deck_json, package_json, out_path)
