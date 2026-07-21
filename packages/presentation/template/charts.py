"""Native, editable PowerPoint charts styled from tokens.

Waterfall and football-field use the standard invisible-base stacked technique so every
bar remains an editable chart series (no pictures, no 3D, no gradients).
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.util import Emu, Pt

from .theme import Theme, set_run_font

CHART_TYPE_MAP = {
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "grouped_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "waterfall": XL_CHART_TYPE.COLUMN_STACKED,       # base series hidden
    "football_field": XL_CHART_TYPE.BAR_STACKED,     # base series hidden
}


def _style_common(chart, theme: Theme, lang: str, *, show_legend: bool) -> None:
    chart.has_title = False
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = theme.muted
    else:
        chart.has_legend = False
    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
        except ValueError:
            continue
        axis.tick_labels.font.size = Pt(9)
        axis.tick_labels.font.color.rgb = theme.muted
        axis.format.line.color.rgb = theme.gridline
        axis.major_tick_mark = XL_TICK_MARK.NONE
        if axis_name == "value_axis":
            axis.has_major_gridlines = True
            axis.major_gridlines.format.line.color.rgb = theme.gridline
            axis.major_gridlines.format.line.width = Pt(0.5)
        else:
            axis.has_major_gridlines = False


def _fill_series(series, color: RGBColor, *, invisible: bool = False, is_line: bool = False) -> None:
    if is_line:
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.25)
        return
    if invisible:
        series.format.fill.background()
        series.format.line.fill.background()
    else:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.fill.background()


def add_chart(slide, theme: Theme, lang: str, *, chart_type: str,
              categories: list[str], series: list[tuple[str, list[float | None]]],
              x: Emu, y: Emu, w: Emu, h: Emu,
              signal_colors: list[RGBColor] | None = None,
              axis_min: float | None = None, axis_max: float | None = None) -> None:
    """series: [(name, values)]. For waterfall/football, series[0] must be the invisible base."""
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    xl_type = CHART_TYPE_MAP[chart_type]
    frame = slide.shapes.add_chart(xl_type, x, y, w, h, data)
    chart = frame.chart
    hidden_base = chart_type in ("waterfall", "football_field")
    multi = len(series) > (2 if hidden_base else 1) and chart_type != "donut"
    if chart_type in ("waterfall", "football_field"):
        multi = False  # colour meaning is stated in the caption; base series stays unlisted
    _style_common(chart, theme, lang, show_legend=multi)
    if axis_min is not None or axis_max is not None:
        axis = chart.value_axis
        if axis_min is not None:
            axis.minimum_scale = axis_min
        if axis_max is not None:
            axis.maximum_scale = axis_max
    is_line = chart_type == "line"
    plot = chart.plots[0]
    plot.gap_width = 60 if not is_line else 150
    if chart_type in ("waterfall", "football_field"):
        plot.gap_width = 45
        try:
            plot.overlap = 100
        except AttributeError:
            pass
    visible_idx = 0
    for i, s in enumerate(chart.series):
        if hidden_base and i == 0:
            _fill_series(s, theme.surface, invisible=True)
            continue
        if signal_colors is not None:
            _fill_series(s, signal_colors[visible_idx % len(signal_colors)], is_line=is_line)
        else:
            _fill_series(s, theme.series_color(visible_idx), is_line=is_line)
        visible_idx += 1
    # data labels only on small single-series charts
    n_points = len(categories)
    if not multi and n_points <= 8 and chart_type not in ("donut", "scatter"):
        target = chart.series[-1] if hidden_base else chart.series[0]
        target.has_data_labels = True
        target.data_labels.font.size = Pt(9)
        target.data_labels.font.color.rgb = theme.body_text
        target.data_labels.number_format = "#,##0.0"
        target.data_labels.number_format_is_linked = False


def add_chart_caption(slide, theme: Theme, lang: str, *, text: str,
                      x: Emu, y: Emu, w: Emu) -> None:
    box = slide.shapes.add_textbox(x, y, w, Emu(int(0.22 * 914400)))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run_font(run, theme, lang=lang, size=theme.size_pt("footnote", lang),
                 color=theme.muted)


def waterfall_series(components: list[tuple[str, float]]) -> tuple[list[str], list[tuple[str, list[float | None]]]]:
    """Build (categories, [base, decrease, increase]) stacked series from signed components.

    First and last components are treated as absolute totals (base pillars).
    """
    cats = [name for name, _ in components]
    base: list[float | None] = []
    total: list[float | None] = []
    inc: list[float | None] = []
    dec: list[float | None] = []
    running = 0.0
    for i, (_name, value) in enumerate(components):
        is_total = i == 0 or i == len(components) - 1
        if is_total:
            base.append(0.0)
            total.append(value)
            inc.append(0.0)
            dec.append(0.0)
            running = value if i == 0 else running
        elif value >= 0:
            base.append(running)
            total.append(0.0)
            inc.append(value)
            dec.append(0.0)
            running += value
        else:
            running += value
            base.append(running)
            total.append(0.0)
            inc.append(0.0)
            dec.append(-value)
    return cats, [("base", base), ("总额/Total", total),
                  ("增加/Increase", inc), ("减少/Decrease", dec)]


def football_series(bars: list[tuple[str, float, float]]) -> tuple[list[str], list[tuple[str, list[float | None]]]]:
    """bars: [(label, low, high)] -> (categories, [invisible base, visible range])."""
    cats = [b[0] for b in bars]
    base = [b[1] for b in bars]
    span = [b[2] - b[1] for b in bars]
    return cats, [("base", base), ("range", span)]
