#!/usr/bin/env python3
"""Build the template + sample decks and run the full QA battery (Issue #24 tests 1-16)."""

from __future__ import annotations

import json
import re
import subprocess
import sys
from pathlib import Path

PRES = Path(__file__).resolve().parents[1]
ROOT = PRES.parents[1]
sys.path.insert(0, str(PRES))

from jsonschema import Draft202012Validator  # noqa: E402
from referencing import Registry, Resource  # noqa: E402

from template.deck import build_from_files  # noqa: E402

FIX = PRES / "fixtures"
OUTD = PRES / "qa"
SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

results: list[tuple[str, bool, str]] = []


def check(name: str, ok: bool, detail: str = "") -> None:
    results.append((name, ok, detail))
    print(("PASS " if ok else "FAIL ") + name + (f"  {detail}" if detail and not ok else ""))


def main() -> int:
    OUTD.mkdir(exist_ok=True)
    # 1. fixtures validate against the existing schemas (read-only use)
    pkg_schema = json.loads((ROOT / "packages/research-schema/research-package.schema.json").read_text())
    deck_schema = json.loads((ROOT / "packages/research-schema/slide-deck.schema.json").read_text())
    registry = Registry().with_resources([
        (pkg_schema["$id"], Resource.from_contents(pkg_schema)),
        (deck_schema["$id"], Resource.from_contents(deck_schema)),
    ])
    pkg = json.loads((FIX / "mock-package.json").read_text())
    errs = list(Draft202012Validator(pkg_schema, registry=registry).iter_errors(pkg))
    check("package-schema", not errs, "; ".join(e.message[:80] for e in errs[:3]))
    decks = {}
    for f in ["deck-sample.zh-CN.json", "deck-sample.en-AU.json", "deck-all-layouts.zh-CN.json"]:
        deck = json.loads((FIX / f).read_text())
        decks[f] = deck
        errs = list(Draft202012Validator(deck_schema, registry=registry).iter_errors(deck))
        check(f"deck-schema:{f}", not errs, "; ".join(e.message[:80] for e in errs[:3]))

    # 2. zh/en reference identical MET/CLM/SRC/ASM ids
    ids = {}
    for f in ["deck-sample.zh-CN.json", "deck-sample.en-AU.json"]:
        ids[f] = sorted(set(re.findall(r"(?:MET|CLM|SRC|ASM)-\d{3}", json.dumps(decks[f]))))
    check("bilingual-id-parity", ids["deck-sample.zh-CN.json"] == ids["deck-sample.en-AU.json"])

    # cross-reference resolution
    known = set(re.findall(r"(?:MET|CLM|SRC|ASM)-\d{3}", json.dumps(pkg)))
    for f, deck in decks.items():
        used = set(re.findall(r"(?:MET|CLM|SRC|ASM)-\d{3}", json.dumps(deck)))
        check(f"refs-resolve:{f}", used <= known, str(sorted(used - known)))

    # 3. layout coverage across the template deck
    layouts = {s["layout"] for s in decks["deck-all-layouts.zh-CN.json"]["slides"]}
    expected = {f"L{i:02d}" for i in range(1, 12)}
    covered = {layout.split("_")[0] for layout in layouts}
    check("layout-coverage-L01-L11", covered >= expected, str(sorted(expected - covered)))

    # build the three PPTX files
    reports = []
    for f, out_name in [("deck-all-layouts.zh-CN.json", "jacaranda-template.pptx"),
                        ("deck-sample.zh-CN.json", "sample-report.zh-CN.pptx"),
                        ("deck-sample.en-AU.json", "sample-report.en-AU.pptx")]:
        report = build_from_files(FIX / f, FIX / "mock-package.json", OUTD / out_name)
        reports.append(report)
        check(f"build+geometry:{out_name}", report["status"] == "pass",
              json.dumps(report["issues"][:2], ensure_ascii=False))
        check(f"16x9:{out_name}", report["page_size_in"] == [13.333, 7.5],
              str(report["page_size_in"]))

    (OUTD / "overflow-report.json").write_text(
        json.dumps(reports, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    # 4. reopen: editable objects, charts are native chart parts, logo present
    from pptx import Presentation
    for out_name in ["jacaranda-template.pptx", "sample-report.zh-CN.pptx",
                     "sample-report.en-AU.pptx"]:
        prs = Presentation(str(OUTD / out_name))
        n_charts = sum(1 for s in prs.slides for sh in s.shapes if sh.has_chart)
        n_tables = sum(1 for s in prs.slides for sh in s.shapes if sh.has_table)
        n_pics = [sum(1 for sh in s.shapes if sh.shape_type == 13) for s in prs.slides]
        check(f"reopen:{out_name}", len(prs.slides) >= 5)
        check(f"native-charts:{out_name}", n_charts >= 4, f"charts={n_charts}")
        check(f"native-tables:{out_name}", n_tables >= 1, f"tables={n_tables}")
        # logo: cover has 1 picture; every content slide exactly 1 (the shield)
        content_ok = all(n == 1 for n in n_pics[1:] )
        check(f"logo-on-slides:{out_name}", n_pics[0] == 1 and content_ok, str(n_pics))

    # 5. missing values render as 暂无数据 / N/A
    zh_text = _all_text(OUTD / "sample-report.zh-CN.pptx")
    en_text = _all_text(OUTD / "sample-report.en-AU.pptx")
    check("missing-zh", "暂无数据" in zh_text)
    check("missing-en", "N/A" in en_text)
    check("disclaimer-zh", "不构成任何投资建议" in zh_text)
    check("disclaimer-en", "not investment advice" in en_text)

    # 6. render previews for visual inspection
    render_ok = True
    for out_name in ["jacaranda-template.pptx", "sample-report.zh-CN.pptx",
                     "sample-report.en-AU.pptx"]:
        sub = OUTD / "previews" / out_name.replace(".pptx", "")
        sub.mkdir(parents=True, exist_ok=True)
        proc = subprocess.run(
            [SOFFICE, "--headless", "--convert-to", "png", "--outdir", str(sub),
             str(OUTD / out_name)], capture_output=True, text=True, timeout=300)
        pdf = subprocess.run(
            [SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", str(sub),
             str(OUTD / out_name)], capture_output=True, text=True, timeout=300)
        render_ok = render_ok and pdf.returncode == 0
    check("preview-render", render_ok)

    ok = all(r[1] for r in results)
    print(f"\n{'ALL PASS' if ok else 'FAILURES PRESENT'} ({sum(1 for r in results if r[1])}/{len(results)})")
    return 0 if ok else 1


def _all_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks)


if __name__ == "__main__":
    sys.exit(main())
