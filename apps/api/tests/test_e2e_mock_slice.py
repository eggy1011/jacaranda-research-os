from __future__ import annotations

import asyncio
import copy
import importlib
import json
import runpy
import socket
import sys
from dataclasses import replace
from pathlib import Path
from typing import Any, cast

import pytest
from jsonschema import ValidationError as JsonSchemaValidationError
from pptx import Presentation
from pydantic import ValidationError

from jacaranda_api.e2e import MockResearchOrchestrator
from jacaranda_api.e2e.cli import main, repository_root
from jacaranda_api.e2e.mock_providers import FixtureAkshareClient, ScriptedMockLLMProvider
from jacaranda_api.e2e.models import DemoRequest, PipelineConfigurationError
from jacaranda_api.e2e.orchestrator import run_pipeline
from jacaranda_api.e2e.presentation import PresentationFailure, TemplatePresentationProvider
from jacaranda_api.e2e.validation import (
    SemanticValidationError,
    load_json,
    validate_decks,
    validate_package,
)
from jacaranda_api.llm.errors import LLMProviderError, RetryExhaustedError
from jacaranda_api.llm.models import LLMResult
from jacaranda_api.market_data.errors import SymbolNormalizationError
from jacaranda_api.market_data.symbols import normalize_symbol

ROOT = Path(__file__).resolve().parents[3]


def read(path: Path) -> dict[str, Any]:
    return cast(dict[str, Any], json.loads(path.read_text(encoding="utf-8")))


class SocketInspectingLLM(ScriptedMockLLMProvider):
    def __init__(self, repository_root: Path) -> None:
        super().__init__(repository_root)
        self._translation_changed = False

    async def run(self, *args: Any, **kwargs: Any) -> LLMResult:
        with pytest.raises(RuntimeError, match="network disabled for mock pipeline"):
            socket.socket()
        result = await super().run(*args, **kwargs)
        if args[0] == "translation" and not self._translation_changed:
            output = copy.deepcopy(result.output)
            output["texts"][0]["en_AU"] = "Deliberately changed translation content."
            self._translation_changed = True
            return result.model_copy(update={"output": output})
        return result


def test_full_socket_blocked_vertical_slice(tmp_path: Path) -> None:
    provider = SocketInspectingLLM(ROOT)
    artifacts = asyncio.run(
        MockResearchOrchestrator(ROOT, llm=provider).run(DemoRequest(), tmp_path / "run")
    )
    package = read(artifacts.research_package)
    assert package["status"] == "verified"
    assert package["company"]["is_mock"] is True
    assert "Deliberately changed translation content." in json.dumps(package)
    assert package["metrics"][16]["metric_id"] == "MET-014"
    assert package["metrics"][16]["value"] == 28.4
    checkpoints = read(artifacts.checkpoints)["checkpoints"]
    assert [item["stage"] for item in checkpoints[:8]] == [
        "S1",
        "S2",
        "S3a",
        "S3b",
        "S3c",
        "S3d",
        "S4",
        "S5",
    ]
    assert all(item["prompt_version"] == "0.1.0" for item in checkpoints)
    assert all(item["attempt_count"] == len(item["attempts"]) for item in checkpoints)
    s6_calls = [data for task, data in provider.calls if task == "translation"]
    assert len(s6_calls) > 1
    assert max(len(call["texts"]) for call in s6_calls) <= 20
    s7_names = [task for task, _ in provider.calls if task.startswith("slide_compression")]
    assert s7_names[0] == "slide_compression_plan"
    assert s7_names[1:15] == ["slide_compression_slide"] * 14
    assert s7_names[15] == "slide_compression_plan"
    assert s7_names[16:] == ["slide_compression_slide"] * 14
    for edition in ("zh-CN", "en-AU"):
        deck = read(artifacts.deck_json[edition])
        report = read(artifacts.overflow_reports[edition])
        assert deck["edition"] == edition
        assert report["status"] == "pass" and report["issues"] == []
        pptx = Presentation(str(artifacts.pptx[edition]))
        assert len(pptx.slides) == 14
        assert any(shape.has_text_frame for shape in pptx.slides[0].shapes)
    manifest = read(artifacts.manifest)
    assert manifest["network"] == "socket-blocked"
    assert all(not Path(item["path"]).is_absolute() for item in manifest["artifacts"])
    assert {path.name for path in artifacts.deck_json.values()} == {
        "slide-deck.zh-CN.json",
        "slide-deck.en-AU.json",
    }
    assert {path.name for path in artifacts.pptx.values()} == {
        "report.zh-CN.pptx",
        "report.en-AU.pptx",
    }
    with pytest.raises(SymbolNormalizationError):
        normalize_symbol("600XXX")


def test_retry_is_local_and_exhaustion_is_typed(tmp_path: Path) -> None:
    retrying = ScriptedMockLLMProvider(ROOT, {"financial_analysis": 1})
    artifacts = asyncio.run(
        MockResearchOrchestrator(ROOT, llm=retrying).run(DemoRequest(), tmp_path / "retry")
    )
    checkpoint = next(
        item
        for item in read(artifacts.checkpoints)["checkpoints"]
        if item["task_name"] == "financial_analysis"
    )
    assert [attempt["status"] for attempt in checkpoint["attempts"]] == [
        "retryable_failed",
        "succeeded",
    ]
    failing = ScriptedMockLLMProvider(ROOT, {"extraction": 3})
    with pytest.raises(RetryExhaustedError):
        asyncio.run(
            MockResearchOrchestrator(ROOT, llm=failing).run(DemoRequest(), tmp_path / "exhausted")
        )


class NonRetryableLLM(ScriptedMockLLMProvider):
    async def run(self, *args: Any, **kwargs: Any) -> Any:
        raise LLMProviderError(code="stop", retryable=False, message="safe stop")


class InvalidOutputLLM(ScriptedMockLLMProvider):
    async def run(self, *args: Any, **kwargs: Any) -> LLMResult:
        result = await super().run(*args, **kwargs)
        return result.model_copy(update={"output": {}})


def _translation_batch(
    scheduler: MockResearchOrchestrator, package: dict[str, Any]
) -> dict[str, Any]:
    return {
        "authoritative_language": "zh_CN",
        "texts": copy.deepcopy(scheduler._localized_texts(package)),
        "translation_flags": [],
        "glossary_flags": [],
    }


def test_non_retryable_failure_stops(tmp_path: Path) -> None:
    with pytest.raises(LLMProviderError, match="safe stop"):
        asyncio.run(
            MockResearchOrchestrator(ROOT, llm=NonRetryableLLM(ROOT)).run(
                DemoRequest(), tmp_path / "stop"
            )
        )


def test_scheduler_defensive_contracts() -> None:
    scheduler = MockResearchOrchestrator(ROOT)
    scheduler._tasks_by_stage["S1"] = ()
    with pytest.raises(ValueError, match="exactly one"):
        scheduler._one_task("S1")
    scheduler = MockResearchOrchestrator(ROOT)
    with pytest.raises(LLMProviderError, match="scheduler rejected"):
        asyncio.run(scheduler._execute("slide_compression_slide", {}))
    with pytest.raises(RetryExhaustedError):
        asyncio.run(
            MockResearchOrchestrator(ROOT, llm=InvalidOutputLLM(ROOT), max_attempts=1)._execute(
                "extraction", {}
            )
        )
    with pytest.raises(AssertionError, match="exhaustive"):
        asyncio.run(MockResearchOrchestrator(ROOT, max_attempts=0)._execute("extraction", {}))
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    task = scheduler._catalog.resolve("translation")

    class BadCatalog:
        def resolve(self, task_name: str) -> Any:
            assert task_name == "translation"
            return replace(task, batching={"max_texts_per_call": "twenty"})

    scheduler._catalog = BadCatalog()  # type: ignore[assignment]
    with pytest.raises(ValueError, match="batch size"):
        asyncio.run(scheduler._translate(package))
    with pytest.raises(ValueError, match="missing path"):
        scheduler._apply_translations(package, [])


def test_translation_output_is_applied_and_only_localised_text_changes() -> None:
    scheduler = MockResearchOrchestrator(ROOT)
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    protected = copy.deepcopy(
        {
            "metrics": package["metrics"],
            "sources": package["sources"],
            "assumptions": package["valuation"]["assumptions"],
            "claim_references": [
                {
                    key: claim.get(key)
                    for key in ("claim_id", "source_ids", "metric_ids")
                }
                for claim in package["claims"]
            ],
        }
    )
    batch = _translation_batch(scheduler, package)
    changed_path = batch["texts"][0]["path"]
    batch["texts"][0]["en_AU"] = "Deliberately changed translation content."

    notes = scheduler._apply_translations(package, [batch])

    assert notes == ()
    assert dict(scheduler._localized_text_targets(package))[changed_path]["en_AU"] == (
        "Deliberately changed translation content."
    )
    assert {
        "metrics": package["metrics"],
        "sources": package["sources"],
        "assumptions": package["valuation"]["assumptions"],
        "claim_references": [
            {key: claim.get(key) for key in ("claim_id", "source_ids", "metric_ids")}
            for claim in package["claims"]
        ],
    } == protected


def test_translation_paths_and_authoritative_content_are_strictly_validated() -> None:
    scheduler = MockResearchOrchestrator(ROOT)
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    unknown = _translation_batch(scheduler, package)
    unknown["texts"][0]["path"] = "company.unknown"
    with pytest.raises(ValueError, match="unknown path"):
        scheduler._apply_translations(package, [unknown])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    duplicate = _translation_batch(scheduler, package)
    duplicate["texts"].append(copy.deepcopy(duplicate["texts"][0]))
    with pytest.raises(ValueError, match="duplicate path"):
        scheduler._apply_translations(package, [duplicate])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    changed_authoritative = _translation_batch(scheduler, package)
    changed_authoritative["texts"][0]["zh_CN"] = "未经授权的修改"
    with pytest.raises(ValueError, match="changed authoritative"):
        scheduler._apply_translations(package, [changed_authoritative])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    changed_number = _translation_batch(scheduler, package)
    number_index = next(
        index for index, text in enumerate(changed_number["texts"]) if "15%" in text["en_AU"]
    )
    changed_number["texts"][number_index]["en_AU"] = "Revenue growth changed."
    with pytest.raises(ValueError, match="protected identifiers or numeric values"):
        scheduler._apply_translations(package, [changed_number])
    with pytest.raises(ValueError, match="text must be a string"):
        scheduler._protected_translation_tokens(1)

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    unsupported_language = _translation_batch(scheduler, package)
    unsupported_language["authoritative_language"] = "fr_FR"
    with pytest.raises(ValueError, match="supported authoritative language"):
        scheduler._apply_translations(package, [unsupported_language])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    invalid_flag = _translation_batch(scheduler, package)
    invalid_flag["translation_flags"] = [
        {"path": "company.unknown", "problem": "not a registered text"}
    ]
    with pytest.raises(ValueError, match="review flag contains an unknown path"):
        scheduler._apply_translations(package, [invalid_flag])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    non_string_flag = _translation_batch(scheduler, package)
    non_string_flag["translation_flags"] = [{"path": 1, "problem": "invalid path type"}]
    with pytest.raises(ValueError, match="flag path must be a string"):
        scheduler._apply_translations(package, [non_string_flag])

    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    en_authoritative = _translation_batch(scheduler, package)
    en_authoritative["authoritative_language"] = "en_AU"
    en_authoritative["texts"][0]["zh_CN"] = "刻意改写的中文翻译内容。"
    scheduler._apply_translations(package, [en_authoritative])
    changed_path = en_authoritative["texts"][0]["path"]
    assert dict(scheduler._localized_text_targets(package))[changed_path]["zh_CN"] == (
        "刻意改写的中文翻译内容。"
    )


def test_translation_and_glossary_flags_are_retained_deterministically() -> None:
    scheduler = MockResearchOrchestrator(ROOT)
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    batch = _translation_batch(scheduler, package)
    midpoint = len(batch["texts"]) // 2
    first = {
        **batch,
        "texts": batch["texts"][:midpoint],
        "translation_flags": [
            {"path": batch["texts"][0]["path"], "problem": "source wording is ambiguous"}
        ],
    }
    second = {
        **batch,
        "texts": batch["texts"][midpoint:],
        "translation_flags": [
            {"path": batch["texts"][1]["path"], "problem": "ambiguous antecedent"},
        ],
        "glossary_flags": [
            {
                "path": batch["texts"][0]["path"],
                "term": "增持",
                "proposed_mapping": "Accumulate",
            }
        ],
    }

    notes = scheduler._apply_translations(package, [first, second])
    scheduler._set_generation_metadata(package, notes)

    assert notes == tuple(sorted(notes))
    assert len(notes) == 3
    assert "translation_flag" in package["generation_metadata"]["notes"]
    assert "glossary_flag" in package["generation_metadata"]["notes"]

    empty_package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    empty_notes = scheduler._apply_translations(
        empty_package, [_translation_batch(scheduler, empty_package)]
    )
    scheduler._set_generation_metadata(empty_package, empty_notes)
    assert "review flags" not in empty_package["generation_metadata"]["notes"]


def test_s7_task_selection_uses_registered_task_names() -> None:
    scheduler = MockResearchOrchestrator(ROOT)
    plan, slide = scheduler._tasks_by_stage["S7"]
    scheduler._tasks_by_stage["S7"] = (replace(plan, input_schema={}), slide)
    assert scheduler._s7_task_names() == (
        "slide_compression_plan",
        "slide_compression_slide",
    )

    scheduler._tasks_by_stage["S7"] = (slide,)
    with pytest.raises(PipelineConfigurationError, match="missing required S7 task") as missing:
        scheduler._s7_task_names()
    assert missing.value.code == "invalid_s7_task_registry"

    scheduler._tasks_by_stage["S7"] = (plan, plan, slide)
    with pytest.raises(PipelineConfigurationError, match="duplicated required S7 task"):
        scheduler._s7_task_names()


def test_mock_market_fixture_lookup_is_identifier_based(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    original_load_json = load_json

    def reordered_fixture(path: Path) -> dict[str, Any]:
        fixture = original_load_json(path)
        fixture["sources"] = list(reversed(fixture["sources"]))
        fixture["metrics"] = list(reversed(fixture["metrics"]))
        return fixture

    monkeypatch.setattr("jacaranda_api.e2e.orchestrator.load_json", reordered_fixture)
    market = asyncio.run(MockResearchOrchestrator(ROOT)._market_data(DemoRequest()))
    assert next(item for item in market["sources"] if item["source_id"] == "SRC-002")[
        "type"
    ] == "market_data_api"
    assert next(item for item in market["metrics"] if item["metric_id"] == "MET-014")[
        "value"
    ] == 28.4

    def missing_metric_fixture(path: Path) -> dict[str, Any]:
        fixture = original_load_json(path)
        fixture["metrics"] = [
            item for item in fixture["metrics"] if item["metric_id"] != "MET-014"
        ]
        return fixture

    monkeypatch.setattr("jacaranda_api.e2e.orchestrator.load_json", missing_metric_fixture)
    with pytest.raises(PipelineConfigurationError, match="exactly one metric_id=MET-014"):
        asyncio.run(MockResearchOrchestrator(ROOT)._market_data(DemoRequest()))

    scheduler = MockResearchOrchestrator(ROOT)
    with pytest.raises(PipelineConfigurationError, match="must provide a metric_id"):
        scheduler._replace_fixture_record({}, "metric_id", "MET-014", {})
    with pytest.raises(PipelineConfigurationError, match="exactly one metric_id=MET-014"):
        scheduler._replace_fixture_record(
            [{"metric_id": "MET-014"}, {"metric_id": "MET-014"}],
            "metric_id",
            "MET-014",
            {},
        )

    class EmptyProviderResult:
        def as_research_fragments(self) -> dict[str, list[dict[str, Any]]]:
            return {"sources": [], "metrics": []}

    async def empty_fetch(*args: Any, **kwargs: Any) -> EmptyProviderResult:
        return EmptyProviderResult()

    monkeypatch.setattr(
        "jacaranda_api.e2e.orchestrator.AkshareMarketDataProvider.fetch", empty_fetch
    )
    with pytest.raises(PipelineConfigurationError, match="exactly one source and one metric"):
        asyncio.run(MockResearchOrchestrator(ROOT)._market_data(DemoRequest()))


def test_mock_boundaries_and_output_safety(tmp_path: Path) -> None:
    with pytest.raises(ValidationError):
        DemoRequest(symbol="NOT-MOCK")  # type: ignore[arg-type]
    with pytest.raises(ValidationError):
        DemoRequest(editions=("en-AU", "zh-CN"))
    with pytest.raises(ValidationError):
        DemoRequest(company_name={"zh_CN": "真实公司", "en_AU": "Real Company"})
    with pytest.raises(ValidationError):
        DemoRequest(as_of_date="2026-07-11")  # type: ignore[arg-type]
    with pytest.raises(ValueError, match="fictional sentinel"):
        asyncio.run(FixtureAkshareClient().fetch_quote("NOT-MOCK"))
    occupied = tmp_path / "occupied"
    occupied.mkdir()
    (occupied / "keep").write_text("x")
    with pytest.raises(ValueError, match="new or empty"):
        asyncio.run(MockResearchOrchestrator(ROOT).run(DemoRequest(), occupied))
    link = tmp_path / "link"
    link.symlink_to(tmp_path / "target", target_is_directory=True)
    with pytest.raises(ValueError, match="symlink"):
        asyncio.run(MockResearchOrchestrator(ROOT).run(DemoRequest(), link))
    traversal = tmp_path / "safe" / ".." / "escape"
    with pytest.raises(ValueError, match="new or empty"):
        asyncio.run(MockResearchOrchestrator(ROOT).run(DemoRequest(), traversal))


def test_semantic_validation_failures() -> None:
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    validate_package(ROOT, package)
    cases = []
    approved = json.loads(json.dumps(package))
    approved["status"] = "approved"
    cases.append(approved)
    sections = json.loads(json.dumps(package))
    sections["sections"][-1] = sections["sections"][0]
    cases.append(sections)
    thesis = json.loads(json.dumps(package))
    thesis["sections"][1]["claim_ids"] = ["CLM-004"]
    cases.append(thesis)
    dangling = json.loads(json.dumps(package))
    dangling["sections"][0]["claim_ids"] = ["CLM-999"]
    cases.append(dangling)
    weak_source = json.loads(json.dumps(package))
    weak_source["sources"][0]["reliability_tier"] = "caution"
    cases.append(weak_source)
    failed_qc = json.loads(json.dumps(package))
    failed_qc["quality"]["checks"][0]["result"] = "fail"
    cases.append(failed_qc)
    for item in cases:
        with pytest.raises((SemanticValidationError, JsonSchemaValidationError)):
            validate_package(ROOT, item)

    decks = {
        edition: load_json(ROOT / f"packages/presentation/fixtures/deck-sample.{edition}.json")
        for edition in ("zh-CN", "en-AU")
    }
    validate_decks(ROOT, package, decks)
    mismatch = json.loads(json.dumps(decks))
    mismatch["zh-CN"]["package_id"] = "wrong"
    with pytest.raises(SemanticValidationError, match="identity"):
        validate_decks(ROOT, package, mismatch)
    unresolved = json.loads(json.dumps(decks))
    unresolved["zh-CN"]["slides"][1]["blocks"][0]["kpi_cards"][0]["number"]["metric_id"] = "MET-999"
    with pytest.raises(SemanticValidationError, match="unresolved"):
        validate_decks(ROOT, package, unresolved)
    unresolved_assumption = json.loads(json.dumps(decks))
    unresolved_assumption["zh-CN"]["slides"][11]["blocks"][0]["football_field"]["assumption_lines"][
        0
    ]["refs"][0]["assumption_id"] = "ASM-999"
    with pytest.raises(SemanticValidationError, match="unresolved"):
        validate_decks(ROOT, package, unresolved_assumption)
    parity = json.loads(json.dumps(decks))
    parity_package = json.loads(json.dumps(package))
    extra_metric = json.loads(json.dumps(parity_package["metrics"][0]))
    extra_metric["metric_id"] = "MET-999"
    parity_package["metrics"].append(extra_metric)
    parity["en-AU"]["slides"][1]["blocks"][0]["kpi_cards"][0]["number"]["metric_id"] = "MET-999"
    with pytest.raises(SemanticValidationError, match="parity"):
        validate_decks(ROOT, parity_package, parity)


def test_presentation_failure_and_mock_provider_guards(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    renderer = TemplatePresentationProvider(ROOT)
    deck = load_json(ROOT / "packages/presentation/fixtures/deck-sample.zh-CN.json")
    package = load_json(ROOT / "packages/presentation/fixtures/mock-package.json")
    template_deck = importlib.import_module("template.deck")

    monkeypatch.setattr(
        template_deck, "build_deck", lambda *args: {"status": "fail", "issues": [{}]}
    )
    with pytest.raises(PresentationFailure):
        renderer.render(deck, package, tmp_path / "failed.pptx")
    presentation_path = str(ROOT / "packages/presentation")
    sys.path.insert(0, presentation_path)
    try:
        with pytest.raises(PresentationFailure):
            renderer.render(deck, package, tmp_path / "failed-existing-path.pptx")
    finally:
        sys.path.remove(presentation_path)
    wrong_identity = json.loads(json.dumps(deck))
    wrong_identity["package_id"] = "RPK-WRONG"
    with pytest.raises(PresentationFailure, match="identity"):
        renderer.render(wrong_identity, package, tmp_path / "wrong.pptx")
    mock = ScriptedMockLLMProvider(ROOT)
    with pytest.raises(ValueError, match="full zh-CN"):
        mock._deck("bilingual-summary")
    with pytest.raises(AssertionError, match="unhandled"):
        mock._script("not_registered", {})


def test_cli_and_root_discovery(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch, capsys: pytest.CaptureFixture[str]
) -> None:
    assert repository_root() == ROOT
    monkeypatch.setattr(sys, "argv", ["jacaranda-mock-e2e", "--output-dir", str(tmp_path / "cli")])
    main()
    assert "manifest.json" in capsys.readouterr().out
    second = run_pipeline(ROOT, tmp_path / "direct")
    assert second.research_package.is_file()
    for relative in (
        "research-package.json",
        "slide-deck.zh-CN.json",
        "slide-deck.en-AU.json",
        "overflow-zh-cn.json",
        "overflow-en-au.json",
        "audit/checkpoints.json",
        "manifest.json",
    ):
        assert (tmp_path / "cli" / relative).read_bytes() == (
            tmp_path / "direct" / relative
        ).read_bytes()
    monkeypatch.setattr(Path, "is_file", lambda self: False)
    with pytest.raises(RuntimeError, match="not found"):
        repository_root()


def test_module_entrypoint_invokes_main(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(
        sys,
        "argv",
        ["jacaranda-mock-e2e", "--output-dir", str(tmp_path / "module")],
    )
    runpy.run_module("jacaranda_api.e2e.cli", run_name="__main__")
    assert (tmp_path / "module/manifest.json").is_file()
